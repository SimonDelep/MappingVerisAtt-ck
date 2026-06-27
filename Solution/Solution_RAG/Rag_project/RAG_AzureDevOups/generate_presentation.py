"""Generate the NordTrail Gear RAG project presentation (.pptx)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT = Path(__file__).parent / "presentation_RAG_NordTrail_Gear.pptx"

# NordTrail-inspired palette: forest green + slate
COLOR_TITLE = RGBColor(0x1B, 0x3D, 0x2F)
COLOR_ACCENT = RGBColor(0x2D, 0x6A, 0x4F)
COLOR_BODY = RGBColor(0x33, 0x33, 0x33)
COLOR_MUTED = RGBColor(0x66, 0x66, 0x66)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_title(shape, text: str, size: int = 32):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITLE
    p.alignment = PP_ALIGN.LEFT


def add_bullets(text_frame, items: list[str], level0_size: int = 18):
    text_frame.clear()
    for i, item in enumerate(items):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(level0_size)
        p.font.color.rgb = COLOR_BODY
        p.space_after = Pt(8)


def add_slide_title_content(prs, title: str, bullets: list[str], subtitle: str | None = None):
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    set_title(slide.shapes.title, title, 28)
    body = slide.placeholders[1].text_frame
    if subtitle:
        p = body.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = COLOR_MUTED
        p.space_after = Pt(12)
        for item in bullets:
            para = body.add_paragraph()
            para.text = item
            para.level = 0
            para.font.size = Pt(17)
            para.font.color.rgb = COLOR_BODY
            para.space_after = Pt(6)
    else:
        add_bullets(body, bullets, 17)
    return slide


def add_section_slide(prs, title: str):
    layout = prs.slide_layouts[5]  # Title only (blank with title)
    slide = prs.slides.add_slide(layout)
    title_shape = slide.shapes.title
    set_title(title_shape, title, 36)
    return slide


def add_title_slide(prs):
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "RAG NordTrail Gear"
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(40)
    slide.shapes.title.text_frame.paragraphs[0].font.bold = True
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = COLOR_TITLE

    sub = slide.placeholders[1]
    tf = sub.text_frame
    tf.text = "Retrieval Augmented Generation avec Azure"
    p = tf.paragraphs[0]
    p.font.size = Pt(22)
    p.font.color.rgb = COLOR_ACCENT

    p2 = tf.add_paragraph()
    p2.text = "Cours IA — UQAC"
    p2.font.size = Pt(16)
    p2.font.color.rgb = COLOR_MUTED

    p3 = tf.add_paragraph()
    p3.text = "Mai 2026"
    p3.font.size = Pt(14)
    p3.font.color.rgb = COLOR_MUTED
    return slide


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- INTRODUCTION ---
    add_title_slide(prs)

    add_slide_title_content(
        prs,
        "Introduction — Contexte",
        [
            "NordTrail Gear : entreprise fictive d'équipement outdoor (randonnée, trail).",
            "Besoin : répondre automatiquement aux questions clients (livraison, tailles, SAV, catalogue).",
            "Solution : RAG (Retrieval Augmented Generation) — le LLM s'appuie sur nos documents internes.",
            "Stack cloud : Azure OpenAI (embeddings + chat) et Azure AI Search (index vectoriel).",
            "Stack locale : ChromaDB pour prototyper l'ingestion sans dépendre du cloud à chaque test.",
        ],
    )

    add_slide_title_content(
        prs,
        "Introduction — Objectifs du projet",
        [
            "Indexer une base documentaire hétérogène (PDF, Markdown, CSV, JSON).",
            "Découper les textes en chunks exploitables pour la recherche sémantique.",
            "Stocker des embeddings et retrouver les passages les plus pertinents.",
            "Générer des réponses ancrées dans le contexte récupéré (pas d'hallucination libre).",
            "Évaluer le système avec des questions types (livraison, tailles, colis endommagé, etc.).",
        ],
    )

    # --- WHAT WE DID ---
    add_section_slide(prs, "Ce que nous avons réalisé")

    add_slide_title_content(
        prs,
        "Architecture globale",
        [
            "Dossier documents/ : 7 fichiers indexés + 1 fichier d'évaluation exclu.",
            "config.py : configuration centralisée (.env) et validation des clés API.",
            "Deux pipelines complémentaires : Azure (production) et Chroma (développement local).",
            "Flux Azure : ingest_azure.py → index → rag_test.py (retrieve + LLM).",
            "Flux local : utils.py + ingest.py + embeddings.py + vectorstore.py.",
        ],
    )

    add_slide_title_content(
        prs,
        "Module config.py",
        [
            "Charge les variables depuis .env (clés Azure OpenAI et Azure AI Search).",
            "Paramètres : endpoints, noms de modèles déployés, index, dossier documents.",
            "Chunking configurable : CHUNK_SIZE=1500 caractères, CHUNK_OVERLAP=150.",
            "Fonction validate_config() : arrêt propre si clés manquantes.",
        ],
    )

    add_slide_title_content(
        prs,
        "Module utils.py — Prétraitement",
        [
            "Chargement multi-formats : PDF (pypdf), Markdown, TXT, CSV, JSON.",
            "Nettoyage : suppression caractères parasites, espaces excessifs.",
            "Découpage par mots (chunk_size / overlap) — évite de couper au milieu d'une phrase.",
            "list_document_files() : parcours récursif du dossier documents/.",
        ],
    )

    add_slide_title_content(
        prs,
        "Module embeddings.py",
        [
            "Client Azure OpenAI initialisé depuis config.py.",
            "get_embedding(text) : un vecteur par texte.",
            "get_embeddings(texts) : boucle séquentielle (facile à déboguer en V1).",
            "Le nom du modèle correspond au deployment Azure (ex. mon-embedding).",
            "Dimension typique : 1536 (text-embedding-3-small).",
        ],
    )

    add_slide_title_content(
        prs,
        "Module vectorstore.py — ChromaDB local",
        [
            "Stockage persistant dans ./db/chroma_store.",
            "Collection nordtrail_documents, distance cosine (adaptée aux embeddings OpenAI).",
            "reset_collection() : réindexation propre à chaque ingestion complète.",
            "query_collection() : top-k chunks les plus proches de la question.",
            "Métadonnées par chunk : source, type de fichier, index du chunk.",
        ],
    )

    add_slide_title_content(
        prs,
        "Module ingest.py — Ingestion locale",
        [
            "Pour chaque document autorisé : load → clean → chunk → embed → add Chroma.",
            "Exclusion explicite : emails_clients_test.csv (réservé à l'évaluation).",
            "Résumé en fin de run : nombre de documents et de chunks indexés.",
            "Commande : python ingest.py",
        ],
    )

    add_slide_title_content(
        prs,
        "Module ingest_azure.py — Ingestion cloud",
        [
            "Extracteurs dédiés par extension (.pdf, .md, .csv, .json).",
            "Chunking par caractères avec chevauchement (fenêtre glissante).",
            "Vectorisation Azure OpenAI puis upload_documents vers Azure AI Search.",
            "Chaque chunk indexé préfixe le nom du fichier pour tracer la source.",
            "Commande : python ingest_azure.py",
        ],
    )

    add_slide_title_content(
        prs,
        "Module rag_test.py — Requête RAG",
        [
            "1. Embedding de la question utilisateur.",
            "2. VectorizedQuery sur le champ vector de l'index Azure.",
            "3. Recherche hybride (texte + vecteur) — top résultats comme contexte.",
            "4. Prompt structuré : « Réponds strictement avec le contexte suivant… »",
            "5. Appel au modèle chat (mon-llm-chat) pour la réponse finale.",
            "10 questions de test intégrées (livraison France, tailles trail, SAV, etc.).",
        ],
    )

    add_slide_title_content(
        prs,
        "Corpus documentaire (documents/)",
        [
            "faq_livraison.md — délais et politiques de livraison",
            "guide_tailles.md — conseils tailles chaussures et sacs",
            "procedure_sav_interne.md — procédure service après-vente",
            "conditions_annulation.md — retours et annulations",
            "catalogue_produits.csv — catalogue produits",
            "clients_exemples.json / commandes_exemples.json — données fictives",
            "emails_clients_test.csv — NON indexé (fichier d'évaluation uniquement)",
        ],
    )

    add_slide_title_content(
        prs,
        "Configuration et exécution",
        [
            "Copier .env.example vers .env et renseigner les clés Azure.",
            "pip install -r requirements.txt (ou dépendances listées dans SETUP.md).",
            "Ingestion : python ingest_azure.py",
            "Test RAG : python rag_test.py",
            "Option locale : python ingest.py puis requêtes via Chroma.",
            "test_foundry.py : vérification connexion Microsoft Foundry (.env.test).",
        ],
    )

    add_slide_title_content(
        prs,
        "Différences clés entre les deux pipelines",
        [
            "Azure : chunking par caractères, index cloud scalable, rag_test.py complet.",
            "Local : chunking par mots (utils), ChromaDB, idéal pour itérer rapidement.",
            "Même source documents/ mais stratégies de découpage différentes.",
            "Les deux s'appuient sur Azure OpenAI pour les embeddings (voie locale).",
        ],
    )

    # --- CONCLUSION ---
    add_section_slide(prs, "Conclusion")

    add_slide_title_content(
        prs,
        "Bilan",
        [
            "Système RAG fonctionnel de bout en bout : ingestion → recherche → génération.",
            "Séparation claire des responsabilités (config, utils, embeddings, ingest, rag).",
            "Documentation métier réaliste pour un cas e-commerce / outdoor.",
            "Gestion des secrets via .env — aucune clé API dans le code source.",
            "Base extensible : ajouter un fichier dans documents/ puis relancer l'ingestion.",
        ],
    )

    add_slide_title_content(
        prs,
        "Limites et pistes d'amélioration",
        [
            "Embeddings un par un : passer au batch pour accélérer l'ingestion.",
            "Augmenter top-k (aujourd'hui souvent k=1) pour un contexte plus riche.",
            "Harmoniser chunking caractères vs mots entre pipelines local et Azure.",
            "Interface utilisateur (Streamlit / API) pour démo interactive.",
            "Métriques d'évaluation automatiques sur emails_clients_test.csv.",
        ],
    )

    add_slide_title_content(
        prs,
        "Merci — Questions ?",
        [
            "Projet : RAG_AzureDevOups",
            "Démonstration live : python ingest_azure.py puis python rag_test.py",
            "Documentation : SETUP.md et canvas rag-nordtrail-presentation",
        ],
    )

    return prs


def main():
    prs = build_presentation()
    prs.save(OUTPUT)
    print(f"Presentation saved: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
